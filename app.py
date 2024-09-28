import streamlit as st
import google.generativeai as genai
import os
import PyPDF2 as pdf
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
import json
import io

# Load environment variables
load_dotenv()

# Configure Gemini API
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Function to get response from Gemini API
def get_gemini_response(input_text):
    model = genai.GenerativeModel('gemini-pro')
    try:
        if not input_text.strip():
            st.error("Input text is empty. Cannot generate response.")
            return "{}"  # Return empty JSON

        response = model.generate_content(input_text)
        if response and response.text:
            return response.text
        else:
            st.error("Received an empty response from the model.")
            return "{}"  # Return empty JSON
    except Exception as e:
        st.error(f"Error while getting response from API: {str(e)}")
        return "{}"  # Return empty JSON

# Function to extract text from uploaded PDF file
def input_pdf_text(uploaded_file):
    reader = pdf.PdfReader(uploaded_file)
    text = []
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text.append(page.extract_text() or "")
    return text

# Function to extract text from uploaded Word document
def input_word_text(uploaded_file):
    doc = Document(uploaded_file)
    text = []
    for para in doc.paragraphs:
        text.append(para.text or "")
    return text

# Function to extract text from uploaded PPT file
def input_ppt_text(uploaded_file):
    presentation = Presentation(uploaded_file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text or "")
    return text

# Function to generate DOCX file from content
def generate_docx(generated_content):
    doc = Document()
    doc.add_heading('Generated Content', level=1)

    for content in generated_content:
        doc.add_heading(f'Page {content["Page"]}', level=2)
        doc.add_paragraph(f"**Explanation:**\n{content.get('Explanation', 'No explanation available.')}")
        doc.add_paragraph(f"**Example:**\n{content.get('Example', 'No example available.')}")
        doc.add_paragraph(f"**Test:**\n{content.get('Test', 'No test available.')}")
        doc.add_paragraph(f"**Solution:**\n{content.get('Solution', 'No solution available.')}")
        doc.add_paragraph()  # Add a blank line for spacing

    byte_io = io.BytesIO()
    doc.save(byte_io)
    byte_io.seek(0)
    return byte_io

# Define input prompts for generating content
input_prompts = {
    "Architecture & Design": """
    You are an expert in Architecture & Design. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Arts": """
    You are an expert in Arts. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Business & Economics": """
    You are an expert in Business & Economics. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Education": """
    You are an expert in Education. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Engineering & Technology": """
    You are an expert in Engineering & Technology. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Environmental Studies": """
    You are an expert in Environmental Studies. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Humanities": """
    You are an expert in Humanities. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Law": """
    You are an expert in Law. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Medicine & Health Sciences": """
    You are an expert in Medicine & Health Sciences. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Natural Sciences": """
    You are an expert in Natural Sciences. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Social Sciences": """
    You are an expert in Social Sciences. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Mathematics & Computer Science": """
    You are an expert in Mathematics & Computer Science. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """,
    "Interdisciplinary Studies": """
    You are an expert in Interdisciplinary Studies. Your task is to explain the content on the given page, provide a relevant example, and create a mini test with solutions.
    
    Page Content: {page_content}
    
    I want the response in the following structured format:
    {{"Explanation": "", "Example": "", "Mini Test": "", "Test Solution": "", "Raw Response": ""}}
    """
}

# Streamlit App
st.set_page_config(page_title="Interactify")
page = st.sidebar.radio("Select Page", ["Home", "Ask Me About Your Slide", "History"])

if page == "Home":
    st.title("Interactify")

    # Dropdown for field selection
    field = st.selectbox("Select Field", [
        "Architecture & Design",
        "Arts",
        "Business & Economics",
        "Education",
        "Engineering & Technology",
        "Environmental Studies",
        "Humanities",
        "Law",
        "Medicine & Health Sciences",
        "Natural Sciences",
        "Social Sciences",
        "Mathematics & Computer Science",
        "Interdisciplinary Studies"
    ])

    # File uploader for slides (PDF, Word, PPT, or text) input
    uploaded_file = st.file_uploader("Upload Your Document (PDF, DOCX, PPTX, TXT)...", type=["pdf", "docx", "pptx", "txt"])

    # Text input for page range
    page_range_input = st.text_input("Any specific page ranges (e.g., 4-7):")

    # Initialize session state for history
    if 'history' not in st.session_state:
        st.session_state.history = []

    # Submit button for processing the document
    submit = st.button("Submit")

    if submit:
        if uploaded_file:
            try:
                # Extract text from the uploaded file
                if uploaded_file.type == "application/pdf":
                    document_text = input_pdf_text(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    document_text = input_word_text(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    document_text = input_ppt_text(uploaded_file)
                elif uploaded_file.type == "text/plain":
                    document_text = uploaded_file.read().decode("utf-8").split('\n')
                else:
                    st.error("Unsupported file type!")
                    st.stop()

                # Process page ranges
                page_ranges = []
                if page_range_input:
                    try:
                        ranges = page_range_input.split(',')
                        for r in ranges:
                            start, end = map(int, r.split('-'))
                            page_ranges.append(range(start - 1, end))
                    except ValueError:
                        st.error("Invalid page range format! Use the format 'start-end'.")
                        st.stop()
                else:
                    page_ranges = [range(len(document_text))]

                # Process selected pages
                generated_content = []

                for range_set in page_ranges:
                    for page_num in range_set:
                        if page_num < len(document_text):
                            st.markdown(f"#### Page {page_num + 1}")
                            page_content = document_text[page_num]

                            # Prepare prompt with extracted page text
                            input_prompt = input_prompts.get(field, "")
                            if not input_prompt:
                                st.error(f"No prompt template available for field: {field}")
                                continue
                            input_prompt_filled = input_prompt.format(page_content=page_content)

                            # Get response from Gemini API
                            response = get_gemini_response(input_prompt_filled)

                            try:
                                # Parse response
                                response_json = json.loads(response)

                                # Display and collect Explanation, Example, Test, Solution 
                                explanation = response_json.get("Explanation", "No explanation available.")
                                example = response_json.get("Example", "No example available.")
                                test = response_json.get("Test", "No test available.")
                                solution = response_json.get("Solution", "No test solution available.")

                                st.markdown("**Explanation:**")
                                st.write(explanation)

                                st.markdown("**Example:**")
                                st.write(example)

                                st.markdown("**Test:**")
                                st.write(test)

                                st.markdown("**Solution:**")
                                st.write(solution)

                                # Collect generated content in JSON format
                                page_content_json = {
                                    "Page": page_num + 1,
                                    "Explanation": explanation,
                                    "Example": example,
                                    "Test": test,
                                    "Solution": solution,
                                }
                                generated_content.append(page_content_json)
                            except json.JSONDecodeError:
                                st.error("Failed to decode JSON response from the model.")

                        else:
                            st.warning(f"Page {page_num + 1} is out of range.")

                # Add generated content to history
                st.session_state.history.append(generated_content)

                # Provide option to copy generated content
                
                generated_content_str = json.dumps(generated_content, indent=2)
                st.code(generated_content_str)
                if st.button("Copy to Clipboard"):
                    st.experimental_set_query_params(text=generated_content_str)
                    st.success("Content copied to clipboard!")

                # Provide download option for DOCX
                if generated_content:
                    docx_file = generate_docx(generated_content)
                    st.download_button(
                        label="Download Generated Content as DOCX",
                        data=docx_file,
                        file_name='generated_content.docx',
                        mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                    )
            except Exception as e:
                st.error(f"An error occurred while processing the file: {str(e)}")

elif page == "Ask Me About Your Slide":
    st.title("Ask Me About Your Slide")

    # Text input for asking questions
    question = st.text_input("Ask a question about your slide:")

    if question:
        # Get response from Gemini API for the question
        response = get_gemini_response(question)
        st.markdown("### Answer:")
        st.write(response)

elif page == "History":
    st.title("History")

    # Display history
    if 'history' in st.session_state:
        for index, entry in enumerate(st.session_state.history):
            st.subheader(f"History Entry {index + 1}")
            for content in entry:
                st.markdown(f"**Page {content['Page']}**")
                st.markdown(f"**Explanation:**\n{content['Explanation']}")
                st.markdown(f"**Example:**\n{content['Example']}")
                st.markdown(f"**Test:**\n{content['Test']}")
                st.markdown(f"**Solution:**\n{content['Solution']}")
                st.write("---")
    else:
        st.write("No history available.")
        
st.markdown("""
<style>
[data-testid="stAppViewContainer"] {
    background-image: url("https://images.unsplash.com/photo-1698945746290-a9d1cc575e77");
    background-size: cover;
    background-repeat: no-repeat;
    background-position: center;
}
</style>
""", unsafe_allow_html=True)
